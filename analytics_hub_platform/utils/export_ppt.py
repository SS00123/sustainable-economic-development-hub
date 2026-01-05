"""
PowerPoint Export Utility
Sustainable Economic Development Analytics Hub
Ministry of Economy and Planning

Generates professional PowerPoint presentations with KPI slides.
"""

from datetime import datetime
from io import BytesIO

import pandas as pd

from typing import Any

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    RGBColor = Any  # type: ignore

from analytics_hub_platform.config.branding import BRANDING
from analytics_hub_platform.config.theme import get_theme


def hex_to_rgb(hex_color: str) -> Any:
    """Convert hex color to RGBColor."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


# =============================================================================
# SLIDE BUILDER FUNCTIONS
# =============================================================================


def _add_title_slide(
    prs: "Presentation",
    title: str,
    subtitle: str | None,
    primary_color: "RGBColor",
) -> None:
    """Add title slide to presentation."""
    white = RGBColor(255, 255, 255)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = primary_color
    background.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = white
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1))
        sub_frame = subtitle_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = subtitle
        sub_para.font.size = Pt(24)
        sub_para.font.color.rgb = white
        sub_para.alignment = PP_ALIGN.CENTER

    # Footer with branding
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.333), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_para = footer_frame.paragraphs[0]
    footer_para.text = f"Ministry of Economy and Planning | {datetime.now().strftime('%B %Y')}"
    footer_para.font.size = Pt(14)
    footer_para.font.color.rgb = RGBColor(200, 200, 200)
    footer_para.alignment = PP_ALIGN.CENTER


def _add_executive_summary_slide(
    prs: "Presentation",
    summary_metrics: dict,
    primary_color: "RGBColor",
) -> None:
    """Add executive summary slide with KPI cards."""
    white = RGBColor(255, 255, 255)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = primary_color
    header.line.fill.background()

    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6))
    h_frame = header_text.text_frame
    h_para = h_frame.paragraphs[0]
    h_para.text = "Executive Summary"
    h_para.font.size = Pt(28)
    h_para.font.bold = True
    h_para.font.color.rgb = white

    # KPI Cards
    _add_metric_cards(slide, summary_metrics, primary_color)


def _add_metric_cards(
    slide,
    summary_metrics: dict,
    primary_color: "RGBColor",
) -> None:
    """Add metric cards to a slide."""
    metrics_list = list(summary_metrics.items())[:8]  # Max 8 metrics

    cols = 4
    card_width = Inches(2.8)
    card_height = Inches(1.8)
    start_x = Inches(0.7)
    start_y = Inches(1.8)
    h_gap = Inches(0.3)
    v_gap = Inches(0.3)

    for idx, (metric_name, metric_value) in enumerate(metrics_list):
        row = idx // cols
        col = idx % cols

        x = start_x + (card_width + h_gap) * col
        y = start_y + (card_height + v_gap) * row

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(248, 250, 252)
        card.line.color.rgb = RGBColor(226, 232, 240)

        # Metric name
        name_box = slide.shapes.add_textbox(
            x + Inches(0.15), y + Inches(0.2), card_width - Inches(0.3), Inches(0.5)
        )
        name_frame = name_box.text_frame
        name_para = name_frame.paragraphs[0]
        name_para.text = str(metric_name)[:25]
        name_para.font.size = Pt(11)
        name_para.font.color.rgb = RGBColor(100, 116, 139)

        # Metric value
        value_box = slide.shapes.add_textbox(
            x + Inches(0.15), y + Inches(0.7), card_width - Inches(0.3), Inches(0.8)
        )
        value_frame = value_box.text_frame
        value_para = value_frame.paragraphs[0]
        value_para.text = str(metric_value)
        value_para.font.size = Pt(28)
        value_para.font.bold = True
        value_para.font.color.rgb = primary_color


def _add_data_table_slide(
    prs: "Presentation",
    data: pd.DataFrame,
    primary_color: "RGBColor",
) -> None:
    """Add data table slide to presentation."""
    white = RGBColor(255, 255, 255)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = primary_color
    header.line.fill.background()

    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6))
    h_frame = header_text.text_frame
    h_para = h_frame.paragraphs[0]
    h_para.text = "Data Overview"
    h_para.font.size = Pt(28)
    h_para.font.bold = True
    h_para.font.color.rgb = white

    # Simple table representation
    display_cols = [
        c for c in ["year", "quarter", "region", "sustainability_index"] if c in data.columns
    ][:4]

    if display_cols:
        rows_data = data.head(10)[display_cols].values.tolist()
        n_rows = len(rows_data) + 1  # +1 for header
        n_cols = len(display_cols)

        table = slide.shapes.add_table(
            n_rows, n_cols, Inches(0.5), Inches(1.5), Inches(12), Inches(4)
        ).table

        # Header row
        for col_idx, col_name in enumerate(display_cols):
            cell = table.cell(0, col_idx)
            cell.text = col_name.replace("_", " ").title()
            cell.fill.solid()
            cell.fill.fore_color.rgb = primary_color
            para = cell.text_frame.paragraphs[0]
            para.font.color.rgb = white
            para.font.bold = True
            para.font.size = Pt(12)

        # Data rows
        for row_idx, row_data in enumerate(rows_data):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)[:20]
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(11)


def _add_thank_you_slide(
    prs: "Presentation",
    primary_color: "RGBColor",
) -> None:
    """Add thank you/contact slide to presentation."""
    white = RGBColor(255, 255, 255)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = primary_color
    background.line.fill.background()

    # Thank you text
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1))
    thanks_frame = thanks_box.text_frame
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.text = "Thank You"
    thanks_para.font.size = Pt(48)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = white
    thanks_para.alignment = PP_ALIGN.CENTER

    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1.5))
    contact_frame = contact_box.text_frame
    contact_para = contact_frame.paragraphs[0]
    contact_para.text = BRANDING["author_name"]
    contact_para.font.size = Pt(20)
    contact_para.font.color.rgb = white
    contact_para.alignment = PP_ALIGN.CENTER

    p2 = contact_frame.add_paragraph()
    p2.text = f"{BRANDING['author_email']} | {BRANDING['author_mobile']}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(200, 200, 200)
    p2.alignment = PP_ALIGN.CENTER


# =============================================================================
# MAIN EXPORT FUNCTIONS
# =============================================================================


def generate_ppt_presentation(
    data: pd.DataFrame,
    title: str = "Sustainability Analytics Report",
    subtitle: str | None = None,
    summary_metrics: dict | None = None,
    language: str = "en",
) -> BytesIO:
    """
    Generate a PowerPoint presentation from indicator data.

    Args:
        data: DataFrame containing indicator data
        title: Presentation title
        subtitle: Optional subtitle
        summary_metrics: Optional dict of key metrics to highlight
        language: Language code (en/ar)

    Returns:
        BytesIO buffer containing the PPTX
    """
    if not PPTX_AVAILABLE:
        raise ImportError(
            "python-pptx is required for PPT export. Install with: pip install python-pptx"
        )

    buffer = BytesIO()
    theme = get_theme()

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    primary_color = hex_to_rgb(theme.colors.primary)

    # Slide 1: Title Slide
    _add_title_slide(prs, title, subtitle, primary_color)

    # Slide 2: Executive Summary
    if summary_metrics:
        _add_executive_summary_slide(prs, summary_metrics, primary_color)

    # Slide 3: Data Table (if data provided)
    if not data.empty:
        _add_data_table_slide(prs, data, primary_color)

    # Slide 4: Thank You / Contact
    _add_thank_you_slide(prs, primary_color)

    # Save to buffer
    prs.save(buffer)
    buffer.seek(0)

    return buffer


def generate_kpi_slide(
    kpi_name: str,
    kpi_value: str,
    trend: str | None = None,
    status: str = "neutral",
) -> BytesIO:
    """
    Generate a single KPI slide.

    Args:
        kpi_name: Name of the KPI
        kpi_value: Current value
        trend: Trend description (e.g., "+5% vs last quarter")
        status: green/amber/red/neutral

    Returns:
        BytesIO buffer containing single-slide PPTX
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx required for PPT export")

    buffer = BytesIO()
    theme = get_theme()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Status color
    status_colors = {
        "green": hex_to_rgb(theme.colors.status_green),
        "amber": hex_to_rgb(theme.colors.status_amber),
        "red": hex_to_rgb(theme.colors.status_red),
        "neutral": hex_to_rgb(theme.colors.text_muted),
    }
    status_color = status_colors.get(status, status_colors["neutral"])

    # KPI Name
    name_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1))
    name_frame = name_box.text_frame
    name_para = name_frame.paragraphs[0]
    name_para.text = kpi_name
    name_para.font.size = Pt(24)
    name_para.font.color.rgb = RGBColor(100, 116, 139)
    name_para.alignment = PP_ALIGN.CENTER

    # KPI Value
    value_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1.5))
    value_frame = value_box.text_frame
    value_para = value_frame.paragraphs[0]
    value_para.text = kpi_value
    value_para.font.size = Pt(72)
    value_para.font.bold = True
    value_para.font.color.rgb = status_color
    value_para.alignment = PP_ALIGN.CENTER

    # Trend
    if trend:
        trend_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(0.5))
        trend_frame = trend_box.text_frame
        trend_para = trend_frame.paragraphs[0]
        trend_para.text = trend
        trend_para.font.size = Pt(18)
        trend_para.font.color.rgb = RGBColor(100, 116, 139)
        trend_para.alignment = PP_ALIGN.CENTER

    prs.save(buffer)
    buffer.seek(0)

    return buffer
